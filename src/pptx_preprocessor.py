"""PptxPreprocessor — convert a .pptx deck into a list of :class:`SlideData`.

Two strategies are exposed:

* :func:`convert_pptx` — the **primary** path. Renders each slide via
  LibreOffice headless (``soffice --headless --convert-to pdf``) followed by
  ``pdf2image`` (poppler under the hood) at 120 DPI by default, encodes the
  page images as JPEG (quality 85) by default, and pulls speaker notes from
  the deck with ``python-pptx``. Both ``image_format`` and ``jpeg_quality``
  are configurable; PNG output remains available as an opt-in. Satisfies
  Requirements 1.1, 1.2, 1.3, 1.4, 1.5, 1.6, 14.1, 14.2.
* :func:`load_from_images` — the **fallback** path. Loads pre-exported slide
  images from a user-supplied directory (natural-sorted by filename) and
  pairs them with speaker notes extracted from an optional ``.pptx`` file.
  Satisfies Requirement 1.7 and the `--images-dir` escape hatch described in
  design Error Scenario 1b.

Both paths return ``list[SlideData]`` with ``index`` matching the slide's
position in the deck (Requirement 1.4).
"""

from __future__ import annotations

import base64
import io
import logging
import os
import re
import shutil
import subprocess
import tempfile
from pathlib import Path
from typing import List, Optional, Union

import pptx
import pptx.exc
from pdf2image import convert_from_path
from PIL import Image

from src.models import SlideData

logger = logging.getLogger(__name__)

# ---------------------------------------------------------------------------
# Constants
# ---------------------------------------------------------------------------

DEFAULT_DPI: int = 120
DEFAULT_IMAGE_FORMAT: str = "jpeg"
DEFAULT_JPEG_QUALITY: int = 85
IMAGE_EXTENSIONS: tuple[str, ...] = (".png", ".jpg", ".jpeg")

# JPEG file signature (first 3 bytes are always FF D8 FF).
_JPEG_SIGNATURE: bytes = b"\xFF\xD8\xFF"
_PNG_SIGNATURE: bytes = b"\x89PNG\r\n\x1a\n"

# Common macOS install location — `shutil.which` misses this because the
# LibreOffice installer does not add `soffice` to PATH by default.
_MACOS_SOFFICE_PATH = "/Applications/LibreOffice.app/Contents/MacOS/soffice"

_LIBREOFFICE_INSTALL_INSTRUCTIONS = (
    "LibreOffice (`soffice`) is required to render .pptx slides to images.\n"
    "Install it with one of:\n"
    "  - macOS:   brew install --cask libreoffice\n"
    "  - Linux:   sudo apt install libreoffice  (or your distro equivalent)\n"
    "  - Windows: download from https://www.libreoffice.org/download/\n"
    "Alternatively, export slides to PNG manually (File → Export → PNG) "
    "and pass them via `--images-dir`."
)


# ---------------------------------------------------------------------------
# Exceptions
# ---------------------------------------------------------------------------


class PptxPreprocessorError(Exception):
    """Base class for all preprocessing errors raised by this module."""


class LibreOfficeNotFoundError(PptxPreprocessorError):
    """Raised when the ``soffice`` binary cannot be located."""


class CorruptedPptxError(PptxPreprocessorError):
    """Raised when the .pptx file cannot be opened or rendered."""


# ---------------------------------------------------------------------------
# Internal helpers
# ---------------------------------------------------------------------------


def _find_soffice() -> str:
    """Locate the LibreOffice ``soffice`` binary.

    Raises:
        LibreOfficeNotFoundError: if ``soffice`` is not on PATH and the
            well-known macOS install path does not exist.
    """
    found = shutil.which("soffice")
    if found:
        return found

    if os.path.exists(_MACOS_SOFFICE_PATH):
        return _MACOS_SOFFICE_PATH

    raise LibreOfficeNotFoundError(_LIBREOFFICE_INSTALL_INSTRUCTIONS)


def _extract_notes(pptx_path: Path) -> List[str]:
    """Extract speaker notes from every slide in ``pptx_path``.

    Returns:
        One entry per slide, in deck order. Empty string for slides without
        notes (Requirement 1.3).

    Raises:
        CorruptedPptxError: if the deck cannot be opened (Requirement 1.5).
    """
    try:
        prs = pptx.Presentation(str(pptx_path))
    except pptx.exc.PackageNotFoundError as exc:
        raise CorruptedPptxError(
            f"Cannot open .pptx file at {pptx_path}: {exc}"
        ) from exc
    except Exception as exc:  # python-pptx wraps a few other parse errors
        raise CorruptedPptxError(
            f"Cannot open .pptx file at {pptx_path}: {exc}"
        ) from exc

    notes: List[str] = []
    for slide in prs.slides:
        if slide.has_notes_slide:
            text = slide.notes_slide.notes_text_frame.text or ""
        else:
            text = ""
        notes.append(text)
    return notes


def _render_slides_via_libreoffice(
    pptx_path: Path,
    dpi: int,
    image_format: str = DEFAULT_IMAGE_FORMAT,
    jpeg_quality: int = DEFAULT_JPEG_QUALITY,
) -> List[bytes]:
    """Render every slide of ``pptx_path`` to image bytes via LibreOffice.

    Pipeline: ``.pptx`` → PDF (LibreOffice headless) → page images
    (pdf2image). Each page is encoded to bytes via Pillow.

    Args:
        pptx_path: Path to the ``.pptx`` file.
        dpi: Render resolution passed to ``pdf2image``.
        image_format: ``"jpeg"`` (default) or ``"png"``. JPEG pages are
            encoded via Pillow at ``jpeg_quality`` after flattening any
            RGBA/LA input onto a white background; PNG pages are saved with
            ``optimize=True`` (Requirement 14.2).
        jpeg_quality: JPEG quality (1-95) used when ``image_format == "jpeg"``.

    The temporary directory holding the intermediate PDF is removed on exit
    via :class:`tempfile.TemporaryDirectory` (Requirement for task 4.5).

    Raises:
        LibreOfficeNotFoundError: if ``soffice`` is not available.
        CorruptedPptxError: if LibreOffice fails to render the deck or the
            expected PDF is not produced.
    """
    soffice = _find_soffice()

    with tempfile.TemporaryDirectory(prefix="pptx_preprocessor_") as temp_dir:
        try:
            subprocess.run(
                [
                    soffice,
                    "--headless",
                    "--convert-to",
                    "pdf",
                    "--outdir",
                    temp_dir,
                    str(pptx_path),
                ],
                check=True,
                capture_output=True,
            )
        except subprocess.CalledProcessError as exc:
            stderr = (exc.stderr or b"").decode("utf-8", errors="replace")
            raise CorruptedPptxError(
                f"LibreOffice failed to convert {pptx_path} to PDF: "
                f"{stderr.strip() or exc}"
            ) from exc

        pdf_path = Path(temp_dir) / (pptx_path.stem + ".pdf")
        if not pdf_path.exists():
            # Some LibreOffice builds emit a differently-cased name; fall back
            # to picking the single PDF in the temp dir.
            pdfs = list(Path(temp_dir).glob("*.pdf"))
            if not pdfs:
                raise CorruptedPptxError(
                    f"LibreOffice did not produce a PDF for {pptx_path}"
                )
            pdf_path = pdfs[0]

        pil_images = convert_from_path(
            str(pdf_path), dpi=dpi, thread_count=min(os.cpu_count() or 1, 8)
        )

        encoded_pages: List[bytes] = []
        for pil_image in pil_images:
            buffer = io.BytesIO()
            if image_format == "jpeg":
                # JPEG has no alpha channel — flatten RGBA/LA onto white.
                if pil_image.mode in ("RGBA", "LA"):
                    rgba = (
                        pil_image
                        if pil_image.mode == "RGBA"
                        else pil_image.convert("RGBA")
                    )
                    background = Image.new("RGB", rgba.size, (255, 255, 255))
                    background.paste(rgba, mask=rgba.split()[3])
                    pil_image = background
                elif pil_image.mode != "RGB":
                    pil_image = pil_image.convert("RGB")
                pil_image.save(
                    buffer,
                    format="JPEG",
                    quality=jpeg_quality,
                    optimize=True,
                )
            else:  # png
                pil_image.save(buffer, format="PNG", optimize=True)
            encoded_pages.append(buffer.getvalue())

        return encoded_pages
    # TemporaryDirectory removes the PDF on exit — satisfies 4.5.


def _natural_key(name: str) -> list:
    """Natural-sort key: splits digit runs so ``slide_2`` < ``slide_10``."""
    parts = re.split(r"(\d+)", name)
    return [int(part) if part.isdigit() else part.lower() for part in parts]


def _image_bytes_as(
    raw: bytes,
    image_format: str = "jpeg",
    jpeg_quality: int = 85,
) -> bytes:
    """Return ``raw`` transcoded to *image_format* (``"jpeg"`` or ``"png"``).

    Fast path: if the bytes already match the target format's magic signature,
    return them unchanged.  Otherwise, open with Pillow and re-encode.

    For JPEG output, RGBA/LA images are flattened onto a white background
    before saving at *jpeg_quality*.  For PNG output, ``optimize=True`` is
    used (Requirement 14.2).
    """
    # Fast path — already in the target format.
    if image_format == "png" and raw.startswith(_PNG_SIGNATURE):
        return raw
    if image_format == "jpeg" and raw.startswith(_JPEG_SIGNATURE):
        return raw

    with Image.open(io.BytesIO(raw)) as img:
        buffer = io.BytesIO()
        if image_format == "jpeg":
            # JPEG has no alpha channel — flatten RGBA/LA onto white.
            if img.mode in ("RGBA", "LA"):
                rgba = img if img.mode == "RGBA" else img.convert("RGBA")
                background = Image.new("RGB", rgba.size, (255, 255, 255))
                background.paste(rgba, mask=rgba.split()[3])
                img = background
            elif img.mode != "RGB":
                img = img.convert("RGB")
            img.save(buffer, format="JPEG", quality=jpeg_quality, optimize=True)
        else:  # png
            if img.mode not in ("RGB", "RGBA", "L"):
                img = img.convert("RGB")
            img.save(buffer, format="PNG", optimize=True)
        return buffer.getvalue()


# ---------------------------------------------------------------------------
# Public API
# ---------------------------------------------------------------------------


def convert_pptx(
    file_path: Union[str, os.PathLike],
    dpi: int = DEFAULT_DPI,
    image_format: str = DEFAULT_IMAGE_FORMAT,
    jpeg_quality: int = DEFAULT_JPEG_QUALITY,
) -> List[SlideData]:
    """Convert a .pptx deck into a list of :class:`SlideData` (primary path).

    Args:
        file_path: Path to the .pptx file.
        dpi: Render resolution for the page output. Defaults to 120
            (Requirement 14.1).
        image_format: Encoding for the rendered slide images — ``"jpeg"``
            (default) or ``"png"``. JPEG is preferred to shrink the vision
            payload (Requirement 14.2). When ``"jpeg"``, RGBA/LA images are
            flattened onto a white background before encoding at
            ``jpeg_quality``; when ``"png"``, images are saved with
            ``optimize=True``.
        jpeg_quality: JPEG quality (1-95) used when
            ``image_format == "jpeg"``. Defaults to 85 (Requirement 14.2).

    Returns:
        One :class:`SlideData` per slide, in deck order.

    Raises:
        FileNotFoundError: if ``file_path`` does not exist.
        ValueError: if ``image_format`` is not ``"jpeg"`` or ``"png"``, or if
            ``jpeg_quality`` is outside the 1-95 range.
        LibreOfficeNotFoundError: if ``soffice`` is not available
            (Requirement 1.6).
        CorruptedPptxError: if the deck cannot be opened or rendered
            (Requirement 1.5).
    """
    normalized_format = image_format.lower()
    if normalized_format not in ("jpeg", "png"):
        raise ValueError(
            f"image_format must be 'jpeg' or 'png', got {image_format!r}"
        )
    if not 1 <= jpeg_quality <= 95:
        raise ValueError(
            f"jpeg_quality must be between 1 and 95, got {jpeg_quality}"
        )

    path = Path(file_path)
    if not path.exists():
        raise FileNotFoundError(f"PPTX file not found: {path}")

    page_images = _render_slides_via_libreoffice(
        path,
        dpi,
        image_format=normalized_format,
        jpeg_quality=jpeg_quality,
    )
    notes_list = _extract_notes(path)

    if len(page_images) != len(notes_list):
        logger.warning(
            "Slide-count mismatch for %s: %d rendered pages vs %d notes entries; "
            "pairing by index and padding notes with empty strings where needed.",
            path,
            len(page_images),
            len(notes_list),
        )

    slides: List[SlideData] = []
    for i, image_bytes in enumerate(page_images):
        notes = notes_list[i] if i < len(notes_list) else ""
        img_b64 = base64.b64encode(image_bytes).decode("ascii")
        slides.append(
            SlideData(index=i, image_base64=img_b64, speaker_notes=notes)
        )
    return slides


def load_from_images(
    pptx_path: Optional[Union[str, os.PathLike]],
    images_dir: Union[str, os.PathLike],
    image_format: str = DEFAULT_IMAGE_FORMAT,
    jpeg_quality: int = DEFAULT_JPEG_QUALITY,
) -> List[SlideData]:
    """Load slides from a directory of pre-exported images (fallback path).

    Args:
        pptx_path: Optional path to the original .pptx. When provided,
            speaker notes are extracted from the deck and paired with the
            images by index. Pass ``None`` to skip notes entirely.
        images_dir: Directory containing one image file per slide. Files are
            natural-sorted by filename (e.g. ``slide_2.png`` before
            ``slide_10.png``). Supported extensions: ``.png``, ``.jpg``,
            ``.jpeg`` — images are transcoded to the chosen *image_format*
            when needed.
        image_format: ``"jpeg"`` (default) or ``"png"``. Controls the
            encoding of the resulting slide images (Requirement 14.3).
        jpeg_quality: JPEG quality (1-95) used when
            ``image_format == "jpeg"``. Defaults to 85.

    Returns:
        One :class:`SlideData` per image, in natural-sorted order.

    Raises:
        FileNotFoundError: if ``images_dir`` is missing or contains no
            supported images.
        CorruptedPptxError: if ``pptx_path`` is provided but cannot be opened.
    """
    images_dir_path = Path(images_dir)
    if not images_dir_path.is_dir():
        raise FileNotFoundError(
            f"Images directory not found: {images_dir_path}"
        )

    candidates = [
        entry
        for entry in images_dir_path.iterdir()
        if entry.is_file() and entry.suffix.lower() in IMAGE_EXTENSIONS
    ]
    if not candidates:
        raise FileNotFoundError(
            f"No PNG/JPG images found in {images_dir_path}"
        )

    candidates.sort(key=lambda entry: _natural_key(entry.name))

    notes_list: List[str] = []
    if pptx_path is not None:
        pptx_path_obj = Path(pptx_path)
        if pptx_path_obj.exists():
            notes_list = _extract_notes(pptx_path_obj)
        else:
            logger.warning(
                "PPTX file %s not found; proceeding without speaker notes.",
                pptx_path_obj,
            )

    if notes_list and len(notes_list) != len(candidates):
        logger.warning(
            "Slide-count mismatch between images (%d) and .pptx notes (%d); "
            "pairing by index and padding shorter list.",
            len(candidates),
            len(notes_list),
        )

    slides: List[SlideData] = []
    for i, image_file in enumerate(candidates):
        raw = image_file.read_bytes()
        image_bytes = _image_bytes_as(raw, image_format, jpeg_quality)
        img_b64 = base64.b64encode(image_bytes).decode("ascii")
        notes = notes_list[i] if i < len(notes_list) else ""
        slides.append(
            SlideData(index=i, image_base64=img_b64, speaker_notes=notes)
        )
    return slides


# ---------------------------------------------------------------------------
# Smoke-test CLI
# ---------------------------------------------------------------------------


if __name__ == "__main__":  # pragma: no cover - smoke-test entry point
    import argparse

    parser = argparse.ArgumentParser(
        description="Smoke-test the PPTX preprocessor."
    )
    parser.add_argument("pptx", nargs="?", help="Path to a .pptx file.")
    parser.add_argument(
        "--images-dir",
        help="Directory of pre-exported slide images (fallback path).",
    )
    parser.add_argument(
        "--dpi", type=int, default=DEFAULT_DPI, help="Render DPI."
    )
    args = parser.parse_args()

    logging.basicConfig(level=logging.INFO)

    if args.images_dir:
        result = load_from_images(args.pptx, args.images_dir)
    elif args.pptx:
        result = convert_pptx(args.pptx, dpi=args.dpi)
    else:
        parser.error("Provide a .pptx path or --images-dir.")

    print(f"Loaded {len(result)} slides.")
